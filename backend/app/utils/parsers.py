import io
import logging
from pathlib import Path
from typing import Optional, Dict
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

logger = logging.getLogger(__name__)

class DocumentParser:
    """Parse different document types and extract text"""
    
    @staticmethod
    def parse_pdf(file_path: Path) -> str:
        """Extract text from PDF"""
        try:
            text_content = []
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(f"[Page {page_num}]\n{page_text}")
            
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {e}")
            return ""
    
    @staticmethod
    def parse_docx(file_path: Path) -> str:
        """Extract text from DOCX"""
        try:
            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            
            # Also extract text from tables
            table_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
            
            all_text = paragraphs + table_text
            return "\n\n".join(all_text)
        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {e}")
            return ""
    
    @staticmethod
    def parse_pptx(file_path: Path) -> str:
        """Extract text from PPTX"""
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"[Slide {slide_num}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            return ""
    
    @staticmethod
    def parse_xlsx(file_path: Path) -> str:
        """Extract text from XLSX"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_content = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Convert to string representation
                sheet_text = f"[Sheet: {sheet_name}]\n"
                sheet_text += df.to_string(index=False)
                text_content.append(sheet_text)
            
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Error parsing XLSX {file_path}: {e}")
            return ""
    
    @staticmethod
    def parse_txt(file_path: Path) -> str:
        """Read plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading TXT {file_path}: {e}")
            return ""
    
    @classmethod
    def parse_document(cls, file_path: Path, mime_type: str) -> Dict[str, str]:
        """
        Parse document based on type and return extracted text
        Returns: dict with 'text' and 'error' keys
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            return {"text": "", "error": "File not found"}
        
        try:
            # Determine parser based on mime type or extension
            ext = file_path.suffix.lower()
            
            if 'pdf' in mime_type or ext == '.pdf':
                text = cls.parse_pdf(file_path)
            elif 'wordprocessing' in mime_type or ext in ['.docx', '.doc']:
                text = cls.parse_docx(file_path)
            elif 'presentation' in mime_type or ext in ['.pptx', '.ppt']:
                text = cls.parse_pptx(file_path)
            elif 'spreadsheet' in mime_type or ext in ['.xlsx', '.xls']:
                text = cls.parse_xlsx(file_path)
            elif 'text/plain' in mime_type or ext == '.txt':
                text = cls.parse_txt(file_path)
            else:
                return {"text": "", "error": f"Unsupported file type: {mime_type}"}
            
            if not text or not text.strip():
                return {"text": "", "error": "No text could be extracted"}
            
            return {"text": text, "error": None}
            
        except Exception as e:
            logger.error(f"Error parsing document {file_path}: {e}")
            return {"text": "", "error": str(e)}